from pptx import Presentation
from edatasheets_creator.plugins.powerpointmap import PowerPointMap
from edatasheets_creator.document.jsondatasheet import JsonDataSheet
from edatasheets_creator.document.jsondatasheetschema import JsonDataSheetSchema
from edatasheets_creator.constants import datasheetconstants
import json


class Plugin:
    """
    Powerpoint plugin class that implements datasheet generation from a PPTX
    """

    def __init__(self):
        """
        Class initialization
        """

    def getSlideDetails(self):
        # get the sheet names
        self.slideDetails = self.map.getSlideDetails()

    def getSlideNumbers(self):
        self.slideNumbers = []
        for value in self.slideDetails:
            self.slideNumbers.append(value['slideNumber'])

    def getSlideDatasheetTitle(self):
        self.slideDatasheetTitle = []
        for value in self.slideDetails:
            self.slideDatasheetTitle.append(value['slideDatasheetTitle'])

    def getPowerPointTitle(self):
        self.powerpointTitle = self.map.getPowerPointDescription()

    def process(self, inputFileName, outputFileName, mapFileName=""):
        """

        Main logic method for PowerPoint processing

        Args:
            inputFileName (PosixPath): Input file name
            outputFileName (PosixPath): Output file name
            mapFileName (PosixPath): Map file to guide parser
        """
        self._fileName = inputFileName
        self._outputFileName = outputFileName
        self._mapFileName = mapFileName

        # Extracting information from Powerpoint map
        self.map = PowerPointMap(mapFileName)
        self.getSlideDetails()
        self.getSlideNumbers()
        self.getSlideDatasheetTitle()
        self.getPowerPointTitle()

        # Create presentation class
        prs = Presentation(inputFileName)

        # Process header information of output file
        jds = JsonDataSheet(self._outputFileName)
        datasheetHeader = jds.setMetadata(datasheetconstants.DATASHEET_GENERAL_TITLE, self.powerpointTitle, self._fileName)

        # Process table in datasheet
        textual_data = []
        data_dict = datasheetHeader
        slide_table_dict = {}
        data_dict[datasheetconstants.DATASHEET_TABLES_LOWER] = []
        headers = []
        slide_table = []
        for number in self.slideNumbers:
            headers = []
            slide_table_dict = {}
            slide_table = []
            slide = prs.slides[number - 1]
            for objectify in slide.shapes:
                # Process text
                if objectify.has_text_frame:
                    textual_data.append(objectify.text)
                    print(objectify.text)
                    print('yes')
                # Process table
                if objectify.has_table:
                    for bits in objectify.table.rows[0].cells:
                        intermediate_text = bits.text_frame.text
                        headers.append(intermediate_text .replace("\n", ""))
                    for i in range(1, len(objectify.table.rows)):
                        row_dict = {}
                        for j in range(len(objectify.table.rows[i].cells)):
                            intermediate_text = objectify.table.rows[i].cells[j].text_frame.text
                            row_dict[headers[j]] = intermediate_text
                            textual_data.append(intermediate_text .replace("\n", ""))
                        slide_table.append(row_dict)
                    slide_table_dict[self.slideDatasheetTitle[number - 1]] = slide_table
                    data_dict[datasheetconstants.DATASHEET_TABLES_LOWER].append(slide_table_dict)
        # Write final output to JSON file
        with open(self._outputFileName, "w", encoding="utf8") as output_json:
            json.dump(data_dict, output_json, indent=2, ensure_ascii=False)
        # Write the schema
        schema = JsonDataSheetSchema(self._outputFileName)
        schema.write()
